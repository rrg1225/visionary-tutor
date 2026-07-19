#!/usr/bin/env python3
"""标准版 pptx_generator 集成测试。"""

from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

_ENGINE = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(_ENGINE))

from pptx_generator import generate_pptx_from_json


def _slide_all_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            t = shape.text_frame.text.strip()
            if t:
                parts.append(t)
    return "\n".join(parts)


class TestStandardGenerator(unittest.TestCase):
    def test_quiz_dict_list_renders_stems_not_repr(self) -> None:
        payload = {
            "topic": "CNN",
            "handout": ["卷积层提取局部特征"],
            "quiz": [
                {"question": "卷积核输出尺寸？"},
                {"stem": "MaxPool 作用？"},
            ],
            "mindmap": "graph TD; A-->B",
            "videoScript": ["镜头1"],
            "citations": ["visionary_global_knowledge"],
        }
        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / "test.pptx"
            generate_pptx_from_json(payload, out)
            prs = Presentation(str(out))
            quiz_slide_text = _slide_all_text(prs.slides[4])
            self.assertIn("卷积核输出尺寸", quiz_slide_text)
            self.assertNotIn("{'question'", quiz_slide_text)

    def test_handout_markdown_stripped(self) -> None:
        payload = {
            "topic": "测试",
            "handout": ["**卷积层**通过 `nn.Conv2d` 提取特征"],
        }
        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / "test.pptx"
            generate_pptx_from_json(payload, out)
            prs = Presentation(str(out))
            handout_text = _slide_all_text(prs.slides[3])
            self.assertNotIn("**", handout_text)
            self.assertNotIn("`", handout_text)
            self.assertIn("卷积层", handout_text)


if __name__ == "__main__":
    unittest.main()
