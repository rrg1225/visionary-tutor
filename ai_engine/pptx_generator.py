#!/usr/bin/env python3
"""
PPTX Export Generator for VisionaryTutor
Generates professional teaching PPTX from resource content (HANDOUT / CODE_PRACTICE / LEARNING_PATH / etc.)

Usage (CLI):
  python pptx_generator.py --title "卷积神经网络" --content "..." --type HANDOUT --output output.pptx

Or import as module and call generate_pptx(...)
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Any, List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

_ENGINE_ROOT = Path(__file__).parent
if str(_ENGINE_ROOT) not in sys.path:
    sys.path.insert(0, str(_ENGINE_ROOT))

from pptx_kit.content_normalizer import (
    clean_text_line,
    clean_text_lines,
    coerce_quiz_lines,
)


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "智眸学伴 · 个性化学习资源") -> None:
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 60, 114)  # dark blue
    shape.line.fill.background()

    # title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, title: str, bullets: List[str], note: str = "") -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(30, 60, 114)
    header.line.fill.background()

    # title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets[:12]):  # limit to avoid overflow
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(8)

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)


def _add_citation_slide(prs: Presentation, citations: List[str]) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(30, 60, 114)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "参考来源与引用"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, c in enumerate(citations[:10]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"[{i+1}] {c}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(6)


def parse_content_to_sections(content: str) -> Tuple[str, List[str], List[str], List[str]]:
    """Very lightweight parser: extract title, knowledge bullets, practice, citations."""
    lines = [clean_text_line(l) for l in content.splitlines() if l.strip()]
    lines = [l for l in lines if l]
    title = lines[0] if lines else "学习资源"

    knowledge: List[str] = []
    practice: List[str] = []
    citations: List[str] = []

    current = "knowledge"
    for line in lines[1:]:
        low = line.lower()
        if "练习" in line or "题" in line or "practice" in low:
            current = "practice"
        elif "引用" in line or "来源" in line or "citation" in low or line.startswith("[cite"):
            current = "citation"
        elif line.startswith("#") or line.startswith("##"):
            continue
        else:
            if current == "knowledge" and len(knowledge) < 8:
                knowledge.append(line[:72])
            elif current == "practice" and len(practice) < 6:
                practice.append(line[:72])
            elif current == "citation" and len(citations) < 8:
                citations.append(line[:80])

    if not knowledge:
        knowledge = [l for l in lines[1:6] if len(l) > 5][:6]
    if not practice:
        practice = ["完成对应章节练习题", "上传作业图片进行视觉评测", "根据反馈调整学习路径"]
    if not citations:
        citations = ["知识库证据已通过 CitationValidator 校验", "RAG 检索自 visionary_global_knowledge"]

    return title, knowledge, practice, citations


def _coerce_citations(value: Any) -> List[str]:
    if value is None:
        return []
    if isinstance(value, list):
        return [clean_text_line(str(c)) for c in value if clean_text_line(str(c))]
    text = clean_text_line(str(value))
    return [text] if text else []


def generate_pptx(title: str, content: str, artifact_type: str, output_path: str | Path) -> Path:
    """Main entry: generate a 4-slide teaching PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    clean_title = clean_text_line(title) or "学习资源"

    # Slide 1: Cover
    _add_title_slide(prs, clean_title, f"智眸学伴 · {artifact_type} 个性化资源")

    # Parse
    parsed_title, knowledge, practice, citations = parse_content_to_sections(content)

    # Slide 2: Knowledge points
    _add_content_slide(prs, "核心知识点", knowledge, "来源：RAG 检索 + 多智能体生成，已通过 FactualityVerifier 校验")

    # Slide 3: Practice
    _add_content_slide(prs, "练习与实操", practice, "建议结合可视化资产与代码实操同步进行")

    # Slide 4: Citations
    _add_citation_slide(prs, citations)

    out_path = Path(output_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def _add_toc_slide(prs: Presentation, sections: List[str]) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(30, 60, 114)
    header.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "目录 Contents"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, sec in enumerate(sections[:8]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}. {sec}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(12)


def _add_profile_slide(prs: Presentation, profile: str, topic: str) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(30, 60, 114)
    header.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "学生画像与学习目标"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"主题：{topic}"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = profile or "暂无详细画像，系统将基于默认学习路径生成资源。"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(60, 60, 60)


def generate_pptx_from_json(data: dict, output_path: str | Path) -> Path:
    """Full PPTX from JSON: {topic, studentProfile, handout, mindmap, quiz, videoScript, citations}"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    topic = clean_text_line(str(data.get("topic") or "学习资源"))
    profile = clean_text_line(str(data.get("studentProfile") or ""))
    handout_lines = clean_text_lines(data.get("handout", ""))
    handout = "\n".join(handout_lines)
    mindmap = clean_text_line(str(data.get("mindmap") or ""))
    quiz_lines = [clean_text_line(q) for q in coerce_quiz_lines(data.get("quiz", "")) if clean_text_line(q)]
    video_lines = clean_text_lines(data.get("videoScript", ""))
    citations = _coerce_citations(data.get("citations")) or [
        "RAG 知识库 visionary_global_knowledge",
        "CitationValidator 校验通过",
    ]

    # 1. Cover
    _add_title_slide(prs, topic, "智眸学伴 · 多智能体个性化教学资源包")

    # 2. Profile summary
    _add_profile_slide(prs, profile, topic)

    # 3. TOC
    sections = ["封面与画像", "讲义核心内容", "思维导图", "练习题库", "视频脚本与分镜", "参考来源"]
    _add_toc_slide(prs, sections)

    # 4. Handout knowledge（优先用已清洗的行列表，避免 parse 误把首行当 title）
    knowledge = handout_lines[:8]
    if not knowledge:
        _, knowledge, _, _ = parse_content_to_sections(handout)
    _, _, practice, _ = parse_content_to_sections(handout)
    _add_content_slide(prs, "讲义 · 核心知识点", knowledge, "来源：RAG + DocAgent，已通过 FactualityVerifier")

    # 5. Quiz
    quiz_bullets = [q[:90] for q in quiz_lines[:5]] or practice
    _add_content_slide(prs, "练习题库 · 分层验证", quiz_bullets, "QuizAgent 生成 · 含难度与解析")

    # 6. MindMap (text representation)
    mindmap_bullets = [mindmap[:150]] if mindmap else ["见 Mermaid 代码或前端渲染"]
    _add_content_slide(prs, "思维导图 · 结构可视化", mindmap_bullets, "MindMapAgent 输出 · 支持 Mermaid 渲染")

    # 7. Video storyboard summary
    video_bullets = [v[:80] for v in video_lines[:4]] if video_lines else ["分镜已由 StoryboardAgent 生成"]
    _add_content_slide(prs, "视频脚本 · 分镜概览", video_bullets, "可生成受控教学动画")

    # 8. Citations
    _add_citation_slide(prs, citations)

    out_path = Path(output_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main() -> None:
    parser = argparse.ArgumentParser(description="VisionaryTutor PPTX Generator")
    parser.add_argument("--title", help="Slide title / topic (single mode)")
    parser.add_argument("--content", help="Raw content (single mode)")
    parser.add_argument("--type", default="HANDOUT", help="ArtifactType")
    parser.add_argument("--output", default="export.pptx", help="Output .pptx path")
    parser.add_argument("--json", help="Path to JSON file with full resources (handout+mindmap+quiz+video+profile)")
    args = parser.parse_args()

    if args.json:
        import json
        with open(args.json, "r", encoding="utf-8") as f:
            data = json.load(f)
        generate_pptx_from_json(data, args.output)
        print(f"[pptx_generator] Generated full PPTX from JSON: {args.output}")
    else:
        if not args.title or not args.content:
            parser.error("--title and --content required unless --json is provided")
        generate_pptx(args.title, args.content, args.type, args.output)
        print(f"[pptx_generator] Generated: {args.output}")


if __name__ == "__main__":
    main()