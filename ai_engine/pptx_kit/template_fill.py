#!/usr/bin/env python3
"""
template_fill.py
将 fill_plan 中的占位符替换到 PPTX 模板 Shape 文本中。
"""

from __future__ import annotations

import logging
import re
from copy import deepcopy
from pathlib import Path
from typing import Any, Dict, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.slide import Slide

from .errors import PremiumExportError

logger = logging.getLogger(__name__)

PLACEHOLDER_PATTERN = re.compile(r"\{\{([^}]+)\}\}")
_SHAPE_ID_PATTERN = re.compile(r"^s(\d{2})_sh(\d+)$")

# 可注入图片的占位符类型（Picture / 通用 Object 槽）
_PICTURE_PLACEHOLDER_TYPES = frozenset({
    PP_PLACEHOLDER.PICTURE,
    PP_PLACEHOLDER.OBJECT,
})


def _replace_in_text(text: str, fill_plan: Dict[str, str]) -> str:
    """将文本中所有 {{key}} 替换为 fill_plan 对应值。"""
    if not text:
        return text

    def replacer(match: re.Match[str]) -> str:
        key = match.group(0)  # 含 {{ }}
        inner = match.group(1).strip()
        # 支持 {{topic}} 或 topic 两种 key
        if key in fill_plan:
            return fill_plan[key]
        inner_key = f"{{{{{inner}}}}}"
        if inner_key in fill_plan:
            return fill_plan[inner_key]
        return fill_plan.get(inner, match.group(0))

    return PLACEHOLDER_PATTERN.sub(replacer, text)


def apply_template(
    fill_plan: Dict[str, str],
    template_pptx: Path,
    output_pptx: Path,
) -> Path:
    """
    遍历模板所有 Shape，替换 {{...}} 占位符并保存。
    若模板中无任何占位符被替换，仍输出文件（降级为原模板副本）。
    """
    if not template_pptx.exists():
        raise PremiumExportError("template_pptx_missing", str(template_pptx))

    prs = Presentation(str(template_pptx))
    replaced_count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if not hasattr(shape, "text_frame") or shape.text_frame is None:
                    continue
                tf = shape.text_frame
                for para in tf.paragraphs:
                    if not para.runs:
                        # 无 run 时直接设 text
                        original = para.text or ""
                        new_text = _replace_in_text(original, fill_plan)
                        if new_text != original:
                            para.text = new_text
                            replaced_count += 1
                    else:
                        for run in para.runs:
                            original = run.text or ""
                            new_text = _replace_in_text(original, fill_plan)
                            if new_text != original:
                                run.text = new_text
                                replaced_count += 1
            except Exception:
                # 单个 shape 失败不中断
                continue

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))

    if replaced_count == 0:
        # 无占位符被替换：可能是空模板，记录但不失败
        pass

    return output_pptx


def _find_shape_by_id(slide: Slide, shape_id: str) -> Optional[Any]:
    """
    根据 slots.json 中的 shape_id（如 s05_sh49）定位 Shape。

    shape_id 编码规则与模板 slots 清单一致：
    全局遍历所有 slide.shapes，递增计数器，格式为 s{slide_idx:02d}_sh{counter}。
    """
    match = _SHAPE_ID_PATTERN.match(shape_id.strip())
    if not match:
        logger.error("[template_fill] shape_id 格式无效: %s", shape_id)
        return None

    target_slide_idx = int(match.group(1))
    target_global_idx = int(match.group(2))

    try:
        presentation = slide.part.package.presentation_part.presentation
    except AttributeError as exc:
        logger.error("[template_fill] 无法从 slide 获取 Presentation: %s", exc)
        return None

    global_counter = 0
    for slide_idx, sld in enumerate(presentation.slides):
        for shape in sld.shapes:
            global_counter += 1
            if slide_idx == target_slide_idx and global_counter == target_global_idx:
                return shape

    logger.error(
        "[template_fill] 未找到 shape_id=%s（slide=%s, global=%s）",
        shape_id,
        target_slide_idx,
        target_global_idx,
    )
    return None


def _insert_into_picture_placeholder(shape: Any, image_path: Path) -> bool:
    """向 Picture Placeholder 注入图片。"""
    if not shape.is_placeholder:
        return False
    ph_type = shape.placeholder_format.type
    if ph_type not in _PICTURE_PLACEHOLDER_TYPES:
        return False
    if not hasattr(shape, "insert_picture"):
        return False
    shape.insert_picture(str(image_path))
    return True


def _replace_shape_with_picture(slide: Slide, shape: Any, image_path: Path) -> bool:
    """
    普通形状：保留几何位置，删除原 shape 后在同位置 add_picture。
    """
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    sp_tree = slide.shapes._spTree  # noqa: SLF001 — python-pptx 删除 shape 的标准做法
    sp_tree.remove(shape._element)
    slide.shapes.add_picture(str(image_path), left, top, width, height)
    return True


def apply_image(slide: Slide, shape_id: str, image_path: str) -> bool:
    """
    向指定 slide 的 shape_id 对应形状注入图片。

    策略：
      1. 若为 Picture Placeholder → shape.insert_picture()
      2. 否则 → 读取几何信息，删除原形状，add_picture 回填

    Args:
        slide: 目标幻灯片对象。
        shape_id: 槽位 ID（与 visionary-deck-v1.slots.json 一致，如 ``s05_sh49``）。
        image_path: 本地 PNG/JPG 文件路径。

    Returns:
        True  — 注入成功。
        False — 找不到形状、图片无效或任何异常；**不向上抛异常**。
    """
    try:
        img_path = Path(image_path)
        if not img_path.is_file() or img_path.stat().st_size == 0:
            logger.error("[template_fill] 图片文件不存在或为空: %s", image_path)
            return False

        shape = _find_shape_by_id(slide, shape_id)
        if shape is None:
            return False

        # 策略 1：正规图片占位符
        if shape.is_placeholder and shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if _insert_into_picture_placeholder(shape, img_path):
                logger.info("[template_fill] 已向占位符 %s 注入图片", shape_id)
                return True
            logger.warning(
                "[template_fill] shape_id=%s 是占位符但非 Picture 类型 (type=%s)，尝试替换策略",
                shape_id,
                shape.placeholder_format.type,
            )

        # 策略 2：普通形状（Rectangle / AutoShape 等）→ 原位替换
        _replace_shape_with_picture(slide, shape, img_path)
        logger.info("[template_fill] 已替换形状 %s 为图片", shape_id)
        return True

    except Exception as exc:  # noqa: BLE001 — 导出链路要求吞掉异常，返回 False
        logger.error(
            "[template_fill] apply_image 失败 (shape_id=%s, image=%s): %s",
            shape_id,
            image_path,
            exc,
        )
        return False
