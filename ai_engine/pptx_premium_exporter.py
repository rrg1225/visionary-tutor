#!/usr/bin/env python3
"""
pptx_premium_exporter.py
Premium PPT 导出 CLI 入口。

流程：guard → load_json → normalize → [mindmap 多模态拦截] → build_fill_plan
      → check_plan → apply → [apply_image] → 清理临时 PNG
任一步失败 → stderr + exit(1) → Java fallback
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import tempfile
import traceback
import uuid
from pathlib import Path
from typing import Any, Dict, Optional, Tuple

# 确保 ai_engine 在 path 中
_ENGINE_ROOT = Path(__file__).parent
if str(_ENGINE_ROOT) not in sys.path:
    sys.path.insert(0, str(_ENGINE_ROOT))

from pptx import Presentation

from pptx_kit.check_plan import check_and_trim_plan
from pptx_kit.content_normalizer import PremiumNormalizer, load_capacity_table
from pptx_kit.errors import PremiumExportError
from pptx_kit.fill_plan_builder import build_fill_plan, load_mapping
from pptx_kit.mermaid_renderer import render_mermaid_png
from pptx_kit.session_loader import load_session_json
from pptx_kit.template_fill import apply_image, apply_template
from pptx_kit.template_guard import TemplateTamperedError, verify_template_bundle

# 思维导图页 slide_index（第 6 页，0-based = 5）
MINDMAP_SLIDE_INDEX = 5
MINDMAP_IMAGE_SHAPE_ID_DEFAULT = "s05_sh50"
MINDMAP_RENDER_SUCCESS_TEXT = "结构详情请参考本页思维导图"
MERMAID_RENDER_TIMEOUT = 8

# Mermaid 语法特征：行首图表类型关键字
_MERMAID_HEAD_PATTERN = re.compile(
    r"^\s*(?:"
    r"graph\s+(?:TD|LR|BT|RL|tb|lr|bt|rl)"
    r"|flowchart"
    r"|sequenceDiagram"
    r"|classDiagram"
    r"|stateDiagram(?:-v2)?"
    r"|erDiagram"
    r"|journey"
    r"|gantt"
    r"|pie(?:\s+\w+)?"
    r"|mindmap"
    r"|timeline"
    r"|gitGraph"
    r"|C4Context"
    r")\b",
    re.IGNORECASE | re.MULTILINE,
)

# ```mermaid ... ``` 围栏
_MERMAID_FENCE_PATTERN = re.compile(
    r"```\s*mermaid\s*\n(.*?)```",
    re.DOTALL | re.IGNORECASE,
)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Premium PPTX Exporter")
    parser.add_argument("--json", required=True, help="输入 payload JSON 文件路径")
    parser.add_argument("--template", required=True, help="模板目录路径")
    parser.add_argument("--output", required=True, help="输出 .pptx 文件路径")
    parser.add_argument("--base-name", default="visionary-deck-v1", help="模板基础名")
    return parser.parse_args()


def _extract_mermaid_source(raw: Any) -> str:
    """从 mindmap 字段提取可渲染的 Mermaid 源码（优先解析 ```mermaid 围栏）。"""
    if raw is None:
        return ""
    text = str(raw).strip()
    if not text:
        return ""
    fence_match = _MERMAID_FENCE_PATTERN.search(text)
    if fence_match:
        return fence_match.group(1).strip()
    return text


def is_mermaid_content(raw: Any) -> bool:
    """判断 mindmap 内容是否具备 Mermaid 图表特征。"""
    if raw is None:
        return False
    text = str(raw).strip()
    if not text:
        return False
    if _MERMAID_FENCE_PATTERN.search(text):
        return True
    source = _extract_mermaid_source(text)
    return bool(_MERMAID_HEAD_PATTERN.search(source))


def _resolve_mindmap_image_shape_id(template_dir: Path, base_name: str) -> str:
    """
    从 slots.json 查找 {{mindmap_image}} 槽位；缺失时使用默认 ID。
    模板尚未注入图片槽时，apply_image 会安全返回 False，不影响文本降级。
    """
    slots_path = template_dir / f"{base_name}.slots.json"
    if not slots_path.is_file():
        return MINDMAP_IMAGE_SHAPE_ID_DEFAULT
    try:
        with open(slots_path, "r", encoding="utf-8") as f:
            slots = json.load(f)
        if isinstance(slots, list):
            for slot in slots:
                if not isinstance(slot, dict):
                    continue
                ph = str(slot.get("placeholder", ""))
                if ph in ("{{mindmap_image}}", "{{mindmapImage}}"):
                    sid = slot.get("shape_id")
                    if isinstance(sid, str) and sid.strip():
                        return sid.strip()
    except (OSError, json.JSONDecodeError):
        pass
    return MINDMAP_IMAGE_SHAPE_ID_DEFAULT


def _try_render_mindmap_png(
    payload: Dict[str, Any],
    normalized: Dict[str, Any],
) -> Tuple[Optional[Path], bool]:
    """
    在 normalize 之后、build_fill_plan 之前拦截 mindmap 多模态渲染。

    - 渲染源使用 **原始 payload**（避免 Normalizer 150 字截断破坏语法）
    - 成功：改写 normalized['mindmap'] 为摘要引导语，返回 (png_path, True)
    - 失败或非 Mermaid：不改动 normalized，返回 (None, False)

    本函数不抛异常，不影响 Phase 1-3 文本链路。
    """
    raw_mindmap = payload.get("mindmap")
    if not is_mermaid_content(raw_mindmap):
        print("[INFO] mindmap 非 Mermaid 图表，走纯文本填充", file=sys.stderr)
        return None, False

    mermaid_source = _extract_mermaid_source(raw_mindmap)
    if not mermaid_source:
        print("[WARN] Mermaid 源码为空，降级为文本", file=sys.stderr)
        return None, False

    temp_png = Path(tempfile.gettempdir()) / f"visionary_mindmap_{uuid.uuid4().hex}.png"
    print("[INFO] 检测到 Mermaid 思维导图，开始渲染 PNG...", file=sys.stderr)

    if not render_mermaid_png(mermaid_source, str(temp_png), timeout=MERMAID_RENDER_TIMEOUT):
        print("[WARN] Mermaid 渲染失败，降级为 {{mindmap}} 文本填充", file=sys.stderr)
        return None, False

    # 渲染成功：缩短文本槽内容，后续 check_plan 仍按标量 mindmap 容量校验
    normalized["mindmap"] = MINDMAP_RENDER_SUCCESS_TEXT
    print(f"[INFO] Mermaid 渲染成功: {temp_png}", file=sys.stderr)
    return temp_png, True


def _apply_mindmap_image(
    output_path: Path,
    temp_png: Path,
    image_shape_id: str,
) -> None:
    """在 apply_template 完成后，向思维导图页注入 PNG（失败不抛异常）。"""
    try:
        prs = Presentation(str(output_path))
        if len(prs.slides) <= MINDMAP_SLIDE_INDEX:
            print(
                f"[WARN] 幻灯片不足 {MINDMAP_SLIDE_INDEX + 1} 页，跳过图片注入",
                file=sys.stderr,
            )
            return
        slide = prs.slides[MINDMAP_SLIDE_INDEX]
        if apply_image(slide, image_shape_id, str(temp_png)):
            prs.save(str(output_path))
            print(f"[INFO] 思维导图 PNG 已注入 shape_id={image_shape_id}", file=sys.stderr)
        else:
            print(
                "[WARN] apply_image 失败，PPT 保留文本摘要（结构详情请参考本页思维导图）",
                file=sys.stderr,
            )
    except Exception as exc:  # noqa: BLE001
        print(f"[WARN] 图片注入异常（已忽略）: {exc}", file=sys.stderr)


def _cleanup_temp_png(temp_png: Optional[Path]) -> None:
    """删除 UUID 临时 PNG，防止磁盘泄漏。"""
    if temp_png is None:
        return
    try:
        temp_png.unlink(missing_ok=True)
    except OSError as exc:
        print(f"[WARN] 清理临时 PNG 失败 ({temp_png}): {exc}", file=sys.stderr)


def main() -> None:
    args = parse_args()
    template_dir = Path(args.template)
    base_name = args.base_name
    temp_png: Optional[Path] = None

    try:
        # ① 模板防篡改校验
        print("[INFO] 模板完整性校验...", file=sys.stderr)
        verify_template_bundle(str(template_dir), base_name)

        # ② 安全加载 session JSON
        payload = load_session_json(Path(args.json))
        print(f"[INFO] 加载 payload，字段数: {len(payload)}", file=sys.stderr)

        # ③ Normalizer 硬截断（从 mapping 容量表驱动）
        normalizer = PremiumNormalizer.from_mapping(template_dir, base_name)
        normalized = normalizer.clamp(payload)
        print("[INFO] 内容归一化完成", file=sys.stderr)

        # ③½ mindmap 多模态拦截（normalize 与 build_fill_plan 之间）
        temp_png, mindmap_rendered = _try_render_mindmap_png(payload, normalized)

        # ④ 构建 fill_plan（mindmap 渲染成功时已替换为摘要引导语）
        mapping = load_mapping(template_dir, base_name)
        fill_plan = build_fill_plan(normalized, mapping)
        print(f"[INFO] fill_plan 条目数: {len(fill_plan)}", file=sys.stderr)

        # ⑤ check-plan 二次校验 + 列表感知截断（Phase 1-3 逻辑不变）
        capacity = load_capacity_table(template_dir / f"{base_name}.mapping.json")
        field_mappings = mapping.get("field_mappings", {})
        if not isinstance(field_mappings, dict):
            field_mappings = {}
        fill_plan = check_and_trim_plan(fill_plan, capacity, field_mappings)
        print("[INFO] check-plan 通过", file=sys.stderr)

        # ⑥ apply 填充模板（文本占位符）
        template_pptx = template_dir / f"{base_name}.pptx"
        output_path = Path(args.output)
        apply_template(fill_plan, template_pptx, output_path)

        # ⑥½ 多模态图片注入（渲染成功时；失败则 PPT 仅含文本摘要）
        if mindmap_rendered and temp_png is not None:
            image_shape_id = _resolve_mindmap_image_shape_id(template_dir, base_name)
            _apply_mindmap_image(output_path, temp_png, image_shape_id)

        print(f"[SUCCESS] 已导出: {output_path}", file=sys.stderr)
        sys.exit(0)

    except TemplateTamperedError as e:
        print(f"[FATAL] 模板校验失败: {e}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)
    except PremiumExportError as e:
        print(f"[FATAL] PremiumExportError({e.code}): {e}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"[FATAL] Premium 导出失败: {type(e).__name__}: {e}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)
    finally:
        _cleanup_temp_png(temp_png)


if __name__ == "__main__":
    main()
